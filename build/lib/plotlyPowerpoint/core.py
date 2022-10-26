import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from numerize import numerize 
from scipy.stats import pearsonr
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
import os

#Define functions for table/cell formatting
def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def _set_cell_border(cell, border_color="000000", border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:NoFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

    return cell

#function to set the powerpoint template being used
def setTemplate(fileName):
    from pptx import Presentation
    
    #validate input
    if type(fileName) != str:
        raise Exception("You must input your filename as a string")

    #Load in template for presentation
    try:
        global prs 
        prs = Presentation(fileName)
    except:
        raise Exception("File not found")
        
#function for setting a color palette for charts
def setColors(colors):
    global colorPalette
    colorPalette = colors
    
#master function for creating slides
def createSlides(charts):
    
    #loop through each item in the array
    for z in range(len(charts)):
        
        chartDefinition = charts[z]
        
        #####################
        ### Prepare Data
        #####################

        #get data defined
        temp = chartDefinition['data']

        #filter data if needed
        if 'filters' in chartDefinition:
            filters = []
            for item in chartDefinition['filters']:
                if item["type"] == "int":
                    statement = "temp['" + item["variable"] + "'] " + item["operation"] + " int(" + item["value"] + ")"
                elif item['type'] == 'str':
                    statement = "temp['" + item["variable"] + "'] " + item["operation"] + " '" + item["value"] + "'"
                elif item['type'] == 'date':
                    statement = "temp['" + item["variable"] + "'] " + item["operation"] + " '" + item["value"] + "'"
                elif (item['type'] == 'list') and (item['operation'] == 'in'):
                    statement = "temp['" + item["variable"] + "'].isin(" + str(item["value"]) + ")"
                elif (item['type'] == 'list') and (item['operation'] == 'not in'):
                    statement = "~temp['" + item["variable"] + "'].isin(" + str(item["value"]) + ")"
                filters.append(statement)

            #filter data
            for i in range(len(filters)):
                temp = temp.loc[eval(filters[i]), :]

        #group data by axis and breakdowns
        if chartDefinition['type'] != 'table':
            #assembe list
            groupList = []
            if 'color' in chartDefinition:
                groupList.append(chartDefinition['color'])

            #add axis
            groupList.append(chartDefinition['axis'])

            #add facet if included
            if 'facet' in chartDefinition:
                groupList.append(chartDefinition['facet'])

            #assemble dictionary for aggregation
            metricDict = {}
            for metric in chartDefinition["metrics"]:
                metricDict[metric["name"]] = metric["method"]

            #finally group and summarise data
            temp = temp.groupby(groupList).agg(metricDict).reset_index()


        #####################
        ### Create and Save Chart
        #####################
        
        #set color palette. If pre-set, define it. If not, use default
        try:
            colorPalette
        except NameError:
            mainColors = px.colors.qualitative.Plotly
        else:
            mainColors = colorPalette

        #line chart
        if chartDefinition['type'] == 'line':
            
            #first, figure out if we have multiple metrics. Chart is very different if multiple
            if len(chartDefinition['metrics']) == 1:

                #Determine if we're grouping by color or not
                if 'color' in chartDefinition:  
                    fig = px.line(temp,
                                  x=chartDefinition['axis'],
                                  y=chartDefinition['metrics'][0]['name'],
                                  color_discrete_sequence= mainColors,
                                  color=chartDefinition['color'])
                else:
                    fig = px.line(temp,
                              x=chartDefinition['axis'],
                              y=chartDefinition['metrics'][0]['name'],
                              color_discrete_sequence=mainColors
                                 )

            else: #we have multiple metrics 

                # Create fig
                fig = go.Figure()

                # Add all lines to the chart
                for i in range(len(chartDefinition['metrics'])):
                    fig.add_trace(go.Scatter(x=temp[chartDefinition['axis']],
                                             y=temp[chartDefinition['metrics'][i]['name']],
                                             mode='lines',
                                             name=chartDefinition['metrics'][i]['prettyName'],
                                             line = dict(color=mainColors[i])
                                            )
                                 )


            #change aesthetics
            fig.update_layout({
                'plot_bgcolor': 'rgba(0, 0, 0, 0)',
                'paper_bgcolor': 'rgba(0, 0, 0, 0)',
            })

            ### Handle all options
            if 'options' in chartDefinition:

                ### Grid lines
                if 'horizontal-grid-lines' in chartDefinition['options']:
                    if chartDefinition['options']['horizontal-grid-lines'] == 'true':
                        fig.update_yaxes(showgrid=True, gridwidth=1, gridcolor='#ebebeb')

                if 'vertical-grid-lines' in chartDefinition['options']:
                    if chartDefinition['options']['vertical-grid-lines'] == 'true':
                        fig.update_xaxes(showgrid=True, gridwidth=1, gridcolor='#ebebeb', title="")

                ### X axis ticks rotation
                if 'x-axis-ticks-angle' in chartDefinition['options']:
                    fig.update_xaxes(nticks=temp[chartDefinition['axis']].nunique(), tickangle=chartDefinition['options']['x-axis-ticks-angle'])


            #update legend
            fig.update_layout(legend=dict(
                orientation="h",
                yanchor="bottom",
                xanchor="center",
                x=.5,
                y=-.3,
                title=""
            ))

            #X axis title
            if 'x-axis-title' in chartDefinition:
                fig.update_layout(
                    xaxis_title=chartDefinition['x-axis-title']
                )

            #Y axis title
            if 'y-axis-title' in chartDefinition:
                fig.update_layout(
                    yaxis_title=chartDefinition['y-axis-title']
                )

        #if bar chart
        if chartDefinition['type'] == 'bar':
            
            #First, define whether or not we have 1 or many metrics
            if len(chartDefinition['metrics']) == 1:
                
                #Find proper orientation of bar chart
                if 'options' in chartDefinition:
                    if 'orientation' in chartDefinition['options']:
                        if chartDefinition['options']['orientation'] == 'horizontal':
                            x = temp[chartDefinition['metrics'][0]['name']]
                            y = temp[chartDefinition['axis']]
                            orien='h'
                        else:
                            x = temp[chartDefinition['axis']]
                            y = temp[chartDefinition['metrics'][0]['name']]
                            orien='v'
                    else:
                        x = temp[chartDefinition['axis']]
                        y = temp[chartDefinition['metrics'][0]['name']]
                        orien='v'
                else:
                    x = temp[chartDefinition['axis']]
                    y = temp[chartDefinition['metrics'][0]['name']]
                    orien='v'
                
                #Setup figure, based on if color is set in function
                if 'color' in chartDefinition:
                    fig = px.bar(temp,
                                 x=x,
                                 y=y,
                                 color=chartDefinition['color'],
                                 orientation=orien,
                                 color_discrete_sequence=mainColors
                                )
                else:
                    fig = px.bar(temp,
                                 x=x,
                                 y=y,
                                 color=groupList[0],
                                 orientation=orien,
                                 color_discrete_sequence=mainColors
                                )
  
            else: #multiple metrics
            
                 # Create fig
                fig = go.Figure()

                # Add all bars to chart
                for i in range(len(chartDefinition['metrics'])):

                    #horizontal or vertical for bar chart
                    if 'options' in chartDefinition:
                        if 'orientation' in chartDefinition['options']:
                            if chartDefinition['options']['orientation'] == 'horizontal':
                                x = temp[chartDefinition['metrics'][i]['name']]
                                y = temp[chartDefinition['axis']]
                                orien='h'
                            else:
                                x = temp[chartDefinition['axis']]
                                y = temp[chartDefinition['metrics'][i]['name']]
                                orien='v'
                        else:
                            x = temp[chartDefinition['axis']]
                            y = temp[chartDefinition['metrics'][i]['name']]
                            orien='v'
                    else:
                        x = temp[chartDefinition['axis']]
                        y = temp[chartDefinition['metrics'][i]['name']]
                        orien='v'

                    #add trace to chart    
                    fig.add_trace(
                        go.Bar(
                            x=x,
                            y=y,
                            name=chartDefinition['metrics'][i]['prettyName'],
                            marker_color=mainColors[i],
                            orientation=orien
                        )
                    ) 

            #change aesthetics
            fig.update_layout({
                'plot_bgcolor': 'rgba(0, 0, 0, 0)',
                'paper_bgcolor': 'rgba(0, 0, 0, 0)',
            })
            
            ### Handle Options
            if 'options' in chartDefinition:
                
                #If horizontal, reverse axis
                if 'orientation' in chartDefinition['options']:
                    if chartDefinition['options']['orientation'] == 'horizontal':
                        fig['layout']['yaxis']['autorange'] = "reversed"

#                 #add data labels
#                 if 'datalabels' in chartDefinition['options']:
#                     if chartDefinition['options']['datalabels'] == 'true':
#                         if chartDefinition['label_type'] == 'money':
#                             textFormat = '%{value:$.2s}'
#                         elif chartDefinition['label_type'] == 'percent':
#                             textFormat = '%{value:.1%}'

#                         fig.update_traces(texttemplate=textFormat, textposition='inside', textangle=0)
#                         fig.update_layout(uniformtext_minsize=12)

            

            #update legend
            fig.update_layout(legend=dict(
                orientation="h",
                yanchor="bottom",
                xanchor="center",
                x=.5,
                y=-.3,
                title=""
            ))
            
            
            
        if chartDefinition['type'] == 'facetLine':
            
            #Create Fig
            facets = temp[chartDefinition['facet']].unique().tolist()
            facetSpacing = chartDefinition['options']['facet-spacing'] if 'facet-spacing' in chartDefinition['options'] else 0.1
            
            if chartDefinition['facet-direction'] == 'rows':
                fig = make_subplots(len(facets), 1, vertical_spacing=facetSpacing)
            else:
                fig = make_subplots(1, len(facets), horizontal_spacing=facetSpacing)

            #add traces for all metrics and all facets
            for i in range(len(chartDefinition['metrics'])):
                for facet in facets:

                    #filter data for only current facet
                    temp2 = temp[temp[chartDefinition['facet']] == facet]
                    position = facets.index(facet)

                    #get proper color for line
                    if 'color-grouping' in chartDefinition['options']:
                        if chartDefinition['options']['color-grouping'] == 'facet':
                            lineColor = mainColors[position]
                        else:
                            lineColor = mainColors[i]
                    else:
                        lineColor = mainColors[i]

                    fig.add_trace(
                        go.Scatter(
                            x=temp2[chartDefinition['axis']],
                            y=temp2[chartDefinition['metrics'][i]['name']],
                            mode='lines',
                            name=facet,
                            line = dict(color=lineColor)
                        ), 
                        position + 1 if chartDefinition['facet-direction'] == 'rows' else 1,
                        position + 1 if chartDefinition['facet-direction'] == 'columns' else 1
                    )
              
            
            #change aesthetics
            fig.update_layout({
                'plot_bgcolor': 'rgba(0, 0, 0, 0)',
                'paper_bgcolor': 'rgba(0, 0, 0, 0)',
            })
            
            
            ### Handle all options
            if 'options' in chartDefinition:

                ### Grid lines
                if 'horizontal-grid-lines' in chartDefinition['options']:
                    if chartDefinition['options']['horizontal-grid-lines'] == 'true':
                        fig.update_yaxes(showgrid=True, gridwidth=1, gridcolor='#ebebeb')

                if 'vertical-grid-lines' in chartDefinition['options']:
                    if chartDefinition['options']['vertical-grid-lines'] == 'true':
                        fig.update_xaxes(showgrid=True, gridwidth=1, gridcolor='#ebebeb')


            #update legend
            fig.update_layout(legend=dict(
                orientation="h",
                yanchor="bottom",
                xanchor="center",
                x=.5,
                y=-.3,
                title=""
            ))

            #X axis title
            if 'x-axis-title' in chartDefinition:
                if chartDefinition['facet-direction'] == 'rows':
                    fig.update_xaxes(title_text=chartDefinition['x-axis-title'], row=len(facets), col=1)
                else:
                    for i in range(len(facets)):
                        fig.update_xaxes(title_text=chartDefinition['x-axis-title'], row=1, col=i+1)

            #Y axis title
            if 'y-axis-title' in chartDefinition:
                if chartDefinition['facet-direction'] == 'rows':
                    for i in range(len(facets)):
                        fig.update_yaxes(title_text=chartDefinition['y-axis-title'], row=i+1, col=1)
                else:
                    fig.update_yaxes(title_text=chartDefinition['y-axis-title'], row=1, col=1)
            
            
        #Facet Bar Chart
        if chartDefinition['type'] == 'facetBar':
            
            #Create Fig
            facets = temp[chartDefinition['facet']].unique().tolist()
            facetSpacing = chartDefinition['options']['facet-spacing'] if 'facet-spacing' in chartDefinition['options'] else 0.1
            
            if chartDefinition['facet-direction'] == 'rows':
                fig = make_subplots(len(facets), 1, vertical_spacing=facetSpacing)
            else:
                fig = make_subplots(1, len(facets), horizontal_spacing=facetSpacing)

            #add traces for all metrics and all facets
            for i in range(len(chartDefinition['metrics'])):
                for facet in facets:

                    #filter data for only current facet
                    temp2 = temp[temp[chartDefinition['facet']] == facet]
                    position = facets.index(facet)

                    #get proper color for line
                    if 'color-grouping' in chartDefinition['options']:
                        if chartDefinition['options']['color-grouping'] == 'facet':
                            barColor = mainColors[position]
                        elif chartDefinition['options']['color-grouping'] == 'axis':
                            axisPoints = temp2[chartDefinition['axis']].unique()
                            barColor = mainColors[0:len(axisPoints)]
                        else:
                            barColor = mainColors[i]
                    else:
                        barColor = mainColors[i]

                    fig.add_trace(
                        go.Bar(
                            x=temp2[chartDefinition['axis']],
                            y=temp2[chartDefinition['metrics'][i]['name']],
                            name=facet,
                            marker=dict(color=barColor)
                        ), 
                        position + 1 if chartDefinition['facet-direction'] == 'rows' else 1,
                        position + 1 if chartDefinition['facet-direction'] == 'columns' else 1
                    )

            #change aesthetics
            fig.update_layout({
                'plot_bgcolor': 'rgba(0, 0, 0, 0)',
                'paper_bgcolor': 'rgba(0, 0, 0, 0)',
            })

#             #make facet titles just the value
#             fig.for_each_annotation(lambda a: a.update(text=a.text.split("=")[1]))

#             #add data labels
#             if chartDefinition['label_type'] == 'normal':
#                 fig.update_traces(texttemplate='%{value:.2s}', textposition='outside', textangle=0)
#             elif chartDefinition['label_type'] == 'money':
#                 fig.update_traces(texttemplate='%{value:$.2s}', textposition='inside', textangle=0)

#             #update size and labels
#             fig.update_xaxes(title_text = "Date", tickfont=dict(size=6))
#             fig.update_yaxes(tickfont=dict(size=6))

            #update legend, margins, font size, etc.
            fig.update_layout(
                legend=dict(
                    orientation="h",
                    yanchor="bottom",
                    xanchor="center",
                    x=.5,
                    y=-.3,
                    title=""
                ),
                margin=dict(
                    l=0, r=0, t=40, b=70
                )
            )

        #Filled line chart
        if chartDefinition['type'] == 'filledLine':
            
            #Figure out if there are multiple metrics. If so, throw an error
            if len(chartDefinition['metrics']) == 1:

                #Determine if we're grouping by color or not
                if 'color' in chartDefinition:  
                    fig = px.area(temp,
                                  x=chartDefinition['axis'],
                                  y=chartDefinition['metrics'][0]['name'],
                                  color_discrete_sequence= mainColors,
                                  color=chartDefinition['color'])
                else:
                    fig = px.area(temp,
                              x=chartDefinition['axis'],
                              y=chartDefinition['metrics'][0]['name'],
                              color_discrete_sequence=mainColors
                                 )

            else: #we have multiple metrics 

               raise ValueError('Filled line charts can only have one metric. Please convert your metrics into a variable:value format and break out the line chart by color')


            #change aesthetics
            fig.update_layout({
                'plot_bgcolor': 'rgba(0, 0, 0, 0)',
                'paper_bgcolor': 'rgba(0, 0, 0, 0)',
            })

            ### Handle all options
            if 'options' in chartDefinition:

                ### Grid lines
                if 'horizontal-grid-lines' in chartDefinition['options']:
                    if chartDefinition['options']['horizontal-grid-lines'] == 'true':
                        fig.update_yaxes(showgrid=True, gridwidth=1, gridcolor='#ebebeb')

                if 'vertical-grid-lines' in chartDefinition['options']:
                    if chartDefinition['options']['vertical-grid-lines'] == 'true':
                        fig.update_xaxes(showgrid=True, gridwidth=1, gridcolor='#ebebeb', title="")

                ### X axis ticks rotation
                if 'x-axis-ticks-angle' in chartDefinition['options']:
                    fig.update_xaxes(nticks=temp[chartDefinition['axis']].nunique(), tickangle=chartDefinition['options']['x-axis-ticks-angle'])


            #update legend
            fig.update_layout(legend=dict(
                orientation="h",
                yanchor="bottom",
                xanchor="center",
                x=.5,
                y=-.3,
                title=""
            ))

            #X axis title
            if 'x-axis-title' in chartDefinition:
                fig.update_layout(
                    xaxis_title=chartDefinition['x-axis-title']
                )

            #Y axis title
            if 'y-axis-title' in chartDefinition:
                fig.update_layout(
                    yaxis_title=chartDefinition['y-axis-title']
                )

        #Facet Fill Line
        if chartDefinition['type'] == 'facetFilledLine':
            
            #throw error if there are multiple metrics
            if len(chartDefinition['metrics']) > 1:
                raise ValueError('Filled line charts can only have one metric. Please convert your metrics into a variable:value format and break out the line chart by color')

            #Create Fig
            facets = temp[chartDefinition['facet']].unique().tolist()
            facetSpacing = chartDefinition['options']['facet-spacing'] if 'facet-spacing' in chartDefinition['options'] else 0.1
            
            if chartDefinition['facet-direction'] == 'rows':
                fig = make_subplots(len(facets), 1, vertical_spacing=facetSpacing)
            else:
                fig = make_subplots(1, len(facets), horizontal_spacing=facetSpacing)

            #Add the figure to each subplot
            facetMemory = []
            for facet in facets:
                
                #filter data for only current facet
                temp2 = temp[temp[chartDefinition['facet']] == facet]
                position = facets.index(facet)
                
                #Add figure, based on whether we're breaking down by color                
                if 'color' in chartDefinition:
                    colorOptions = list(temp2[chartDefinition['color']].unique())
                    for clr in colorOptions:
                        
                        #set parameters we need later
                        colorPosition = colorOptions.index(clr)
                        showLegend = False if clr in facetMemory else True
                        
                        #form new temp
                        temp3 = temp2[temp2[chartDefinition['color']] == clr]
                        
                        #add trace
                        fig.add_trace(go.Scatter(
                                x=temp3[chartDefinition['axis']],
                                y=temp3[chartDefinition['metrics'][0]['name']],
                                hoverinfo='x+y',
                                mode='lines',
                                stackgroup='one',
                                fill='tonexty',
                                name=clr,
                                legendgroup=clr,
                                showlegend=showLegend,
                                line=dict(width=0.5, color=mainColors[colorPosition])
                            ),
                            position + 1 if chartDefinition['facet_direction'] == 'rows' else 1,
                            position + 1 if chartDefinition['facet_direction'] == 'columns' else 1
                        )
                        
                        #add memory that we now used this color option within the faceting
                        facetMemory.append(clr)

                        
                else:
                    fig.add_trace(go.Scatter(
                            x=temp2[chartDefinition['axis']],
                            y=temp2[chartDefinition['metrics'][0]['name']],
                            hoverinfo='x+y',
                            mode='lines',
                            fill='tonexty',
                            name=facet,
                            line=dict(width=0.5)
                        ),
                        position + 1 if chartDefinition['facet_direction'] == 'rows' else 1,
                        position + 1 if chartDefinition['facet_direction'] == 'columns' else 1
                    )
            
            #change aesthetics
            fig.update_layout({
                'plot_bgcolor': 'rgba(0, 0, 0, 0)',
                'paper_bgcolor': 'rgba(0, 0, 0, 0)',
            })
            
            
            ### Handle all options
            if 'options' in chartDefinition:

                ### Grid lines
                if 'horizontal-grid-lines' in chartDefinition['options']:
                    if chartDefinition['options']['horizontal-grid-lines'] == 'true':
                        fig.update_yaxes(showgrid=True, gridwidth=1, gridcolor='#ebebeb')

                if 'vertical-grid-lines' in chartDefinition['options']:
                    if chartDefinition['options']['vertical-grid-lines'] == 'true':
                        fig.update_xaxes(showgrid=True, gridwidth=1, gridcolor='#ebebeb')

            #update legend
            fig.update_layout(legend=dict(
                orientation="h",
                yanchor="bottom",
                xanchor="center",
                x=.5,
                y=-.3,
                title=""
            ))

            #X axis title
            if 'x-axis-title' in chartDefinition:
                if chartDefinition['facet_direction'] == 'rows':
                    fig.update_xaxes(title_text=chartDefinition['x-axis-title'], row=len(facets), col=1)
                else:
                    for i in range(len(facets)):
                        fig.update_xaxes(title_text=chartDefinition['x-axis-title'], row=1, col=i+1)

            #Y axis title
            if 'y-axis-title' in chartDefinition:
                if chartDefinition['facet_direction'] == 'rows':
                    for i in range(len(facets)):
                        fig.update_yaxes(title_text=chartDefinition['y-axis-title'], row=i+1, col=1)
                else:
                    fig.update_yaxes(title_text=chartDefinition['y-axis-title'], row=1, col=1)

            
        #Global options to handle outside of individual chart sections
        if 'print-chart' in chartDefinition:
            if chartDefinition['print-chart'] == 'true':
                fig.show()

        #save figure
        if chartDefinition['type'] != 'table':

            #check if the folder for charts exists. If not, create it
            if not os.path.exists('charts'):
                os.makedirs('charts')
            
            #setup params
            filename = 'charts/chart' + str(z) + '.png'
            
            #save out the files
            if chartDefinition['type'] == 'barsubplot':
                fig.write_image(filename, scale=2, width=1.1, height=1)
            elif chartDefinition['name'] == 'Lead Quality - Lead Status Over Time':
                fig.update_layout(margin=dict(r=0))
                fig.write_image(filename, scale=2, width=2, height=1.7)
            else:
                fig.write_image(filename, scale=2)


        #####################
        ### Create Slide and insert image + info
        #####################

        #create slide
        layout = prs.slide_layouts[chartDefinition['item-index']['slide']]
        slide = prs.slides.add_slide(layout)

        #set title and subtitle
        if 'name' in chartDefinition:
            slide.placeholders[chartDefinition['item-index']['title']].text = chartDefinition['name']

        #insert placeholder if desired, otherwise delete
        if "description" in chartDefinition:
            slide.placeholders[chartDefinition['item-index']['description']].text = chartDefinition['description']

        #if we are inserting a plotly image
        if chartDefinition['type'] != 'table':

            #insert image
            picture = slide.placeholders[chartDefinition['item-index']['chart']].insert_picture(filename)

            
        else:
            #insert table
            shape = slide.placeholders[chartDefinition['item-index']['chart']].insert_table(rows=len(temp)+1, cols=len(temp.columns))
            table = shape.table
            
            #iterate through every row and column and place the value that is present in the df
            #for loop for the rows
            for i in range(len(temp) + 1):
                #for each row, get the value of the column
                for i2 in range(len(temp.columns)):
                    cell = table.cell(i,i2)
                    #if we're dealing with the header
                    if i == 0:
                        cell.text = temp.columns[i2]
                    else:
                        text = temp.iloc[i-1, i2]
                        textFormat = chartDefinition['column_formats'][i2]
                        
                        if textFormat == 'number':
                            cell.text = str(int(text))
                        elif textFormat == 'money':
                            cell.text = "$" + str(int(text))
                        elif textFormat == 'percent':
                            cell.text = str(int(text * 100)) + "%"
                        elif textFormat == 'twoDigitNum':
                            cell.text = str(round(text, 2))
                        elif textFormat == 'date':
                            cell.text = text.strftime('%m/%d/%Y')
                        else:
                            cell.text = str(text)
                        
            #central formatting for every cell
            for i in range(len(temp) + 1):
                for i2 in range(len(temp.columns)):
                    #Remove the border for each cell
                    cell = table.cell(i,i2)
                    cell = _set_cell_border(cell)
                    
                    #format color
                    cell = table.cell(i,i2)
                    paragraph = cell.text_frame.paragraphs[0]
                    colorString = chartDefinition['text_color']
                    color = RGBColor.from_string(colorString.replace('#',''))
                    paragraph.font.color.rgb = color

            #If we need to change the header fill
            if 'header_fill_color' in chartDefinition:
                for i in range(len(temp.columns)):
                    cell = table.cell(0,i)
                    cell.fill.solid()
                    colorString = chartDefinition['header_fill_color']
                    color = RGBColor.from_string(colorString.replace('#',''))
                    cell.fill.fore_color.rgb = color
                    
            #If we need to change the header text color
            if 'header_text_color' in chartDefinition:
                for i in range(len(temp.columns)):
                    cell = table.cell(0,i)
                    colorString = chartDefinition['header_text_color']
                    color = RGBColor.from_string(colorString.replace('#',''))
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.color.rgb = color

            #If we need to change the fill color for each cell
            if 'fill_color' in chartDefinition:
                #get the data for the fill coloring
                fillData = chartDefinition['fill_color']
                fillData = fillData.replace("#", '', regex=True)
                
                #loop through each cell
                for i in range(len(temp) + 1):
                    for i2 in range(len(temp.columns)):
                        #skip the header
                        if i != 0:
                            cell = table.cell(i,i2)
                            cell.fill.solid()
                            color = RGBColor.from_string(fillData.iloc[i-1, i2])
                            cell.fill.fore_color.rgb = color

            ### Now center the table in the middle of the slide
            #get base variables
            slideHeight = 5143500
            titleHeight = slide.placeholders[chartDefinition['item-index']['title']].height if 'title' in chartDefinition['item-index'] else 0
            tableHeight = slide.placeholders[chartDefinition['item-index']['chart']].height

            #calculate where the table needs to start
            middleOfSlide = int(slideHeight / 2) + int(titleHeight / 2)
            halfTableHeight = int(tableHeight / 2)
            idealTableStart = int(middleOfSlide - halfTableHeight)

            #set the top of the table
            slide.placeholders[chartDefinition['item-index']['chart']].top = idealTableStart

    #finally save out file
    prs.save("output.pptx")